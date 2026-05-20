"""HTML レポートの全内容を PPTX に変換する。

全 20+ セクションをすべてカバーし、「スキップ推奨」という結論に向かう構成。
工学院大学ロゴ準拠の navy/yellow/white パレットで統一。
スライドサイズ 16:9 (13.333" × 7.5")。

出力: /tmp/school/BVH動作解析_平松.pptx (~55スライド)
"""
from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from PIL import Image

# ===== 色 =====
NAVY = RGBColor(0x2C, 0x41, 0x98)
NAVY_LIGHT = RGBColor(0x42, 0x56, 0xAD)
YELLOW = RGBColor(0xFD, 0xD0, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHTGRAY = RGBColor(0xEE, 0xEE, 0xEE)
DARKGRAY = RGBColor(0x23, 0x18, 0x15)
MUTED_BG = RGBColor(0xF5, 0xF8, 0xFD)

LOGO_PATH = Path("/tmp/school/assets/kogakuin_logo.png")
RG = Path("/tmp/school/reading_guide")
CWT = Path("/tmp/school/analysis/cwt")
DTW = Path("/tmp/school/analysis/dtw")
LG = Path("/tmp/school/legacy_graphs")
OUT = Path("/tmp/school/BVH動作解析_平松.pptx")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height
BLANK = prs.slide_layouts[6]


# ===== ヘルパ =====
def add_logo(slide, scale: float = 1.0):
    w = Inches(1.6 * scale); h = Inches(0.55 * scale)
    left = SLIDE_W - w - Inches(0.3); top = Inches(0.25)
    slide.shapes.add_picture(str(LOGO_PATH), left, top, width=w, height=h)


def add_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = color; bg.line.fill.background()
    spTree = bg._element.getparent()
    spTree.remove(bg._element); spTree.insert(2, bg._element)


def add_textbox(slide, left, top, width, height, text, *,
                size=16, bold=False, color=DARKGRAY,
                align=PP_ALIGN.LEFT, font="Hiragino Sans", line_spacing=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0); tf.margin_right = Inches(0)
    tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]; p.alignment = align
    if line_spacing: p.line_spacing = line_spacing
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color; run.font.name = font
    return tb


def add_bullets(slide, left, top, width, height, items, *,
                size=14, color=DARKGRAY, font="Hiragino Sans"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run(); run.text = f"●  {it}"
        run.font.size = Pt(size); run.font.color.rgb = color
        run.font.name = font
        p.space_after = Pt(5)


def add_section_badge(slide, text):
    w = Inches(2.8); h = Inches(0.45)
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.32), w, h)
    box.fill.solid(); box.fill.fore_color.rgb = NAVY; box.line.fill.background()
    tf = box.text_frame; tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = text
    r.font.size = Pt(13); r.font.bold = True
    r.font.color.rgb = WHITE; r.font.name = "Hiragino Sans"


def add_title(slide, text, top_in=1.0, size=30, color=NAVY):
    add_textbox(slide, Inches(0.5), Inches(top_in), Inches(12.4), Inches(0.8),
                text, size=size, bold=True, color=color)


def add_yellow_underline(slide, top_in, width_in=1.5):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.5), Inches(top_in),
                                  Inches(width_in), Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = YELLOW; bar.line.fill.background()


def add_image_centered(slide, img_path, top_in=1.9,
                       max_width_in=11.5, max_height_in=5.0):
    with Image.open(img_path) as im:
        iw, ih = im.size
    aspect = iw / ih
    if max_width_in / aspect <= max_height_in:
        w = Inches(max_width_in); h = Inches(max_width_in / aspect)
    else:
        h = Inches(max_height_in); w = Inches(max_height_in * aspect)
    left = (SLIDE_W - w) / 2; top = Inches(top_in)
    slide.shapes.add_picture(str(img_path), left, top, width=w, height=h)


def add_finding(slide, label, title, body, top_in=5.6, height_in=1.6):
    left = Inches(0.5); top = Inches(top_in)
    w = Inches(12.333); h = Inches(height_in)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    card.fill.solid(); card.fill.fore_color.rgb = NAVY; card.line.fill.background()
    card.adjustments[0] = 0.04
    # Yellow label
    lb_w = Inches(1.8); lb_h = Inches(0.35)
    lb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.3),
                                top + Inches(0.2), lb_w, lb_h)
    lb.fill.solid(); lb.fill.fore_color.rgb = YELLOW; lb.line.fill.background()
    tf = lb.text_frame; tf.margin_left = Inches(0.1); tf.margin_top = Inches(0.02)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size = Pt(10); r.font.bold = True
    r.font.color.rgb = NAVY; r.font.name = "Hiragino Sans"
    add_textbox(slide, left + Inches(0.3), top + Inches(0.6),
                w - Inches(0.6), Inches(0.45),
                title, size=16, bold=True, color=WHITE)
    add_textbox(slide, left + Inches(0.3), top + Inches(1.0),
                w - Inches(0.6), Inches(height_in - 1.0),
                body, size=11, color=WHITE)


def add_table(slide, headers, rows, left_in, top_in, width_in, height_in,
              header_color=NAVY, font_size=12):
    table = slide.shapes.add_table(
        rows=len(rows) + 1, cols=len(headers),
        left=Inches(left_in), top=Inches(top_in),
        width=Inches(width_in), height=Inches(height_in)).table
    for j, h in enumerate(headers):
        c = table.cell(0, j); c.text = h
        c.fill.solid(); c.fill.fore_color.rgb = header_color
        for p in c.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = WHITE; r.font.bold = True
                r.font.size = Pt(font_size); r.font.name = "Hiragino Sans"
    for i, row in enumerate(rows, start=1):
        for j, v in enumerate(row):
            c = table.cell(i, j); c.text = str(v)
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    r.font.color.rgb = DARKGRAY
                    r.font.size = Pt(font_size - 1)
                    r.font.name = "Hiragino Sans"


# ===== スライド =====
def s_title():
    s = prs.slides.add_slide(BLANK); add_bg(s, WHITE)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.25))
    bar.fill.solid(); bar.fill.fore_color.rgb = YELLOW; bar.line.fill.background()
    nv = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), SLIDE_H)
    nv.fill.solid(); nv.fill.fore_color.rgb = NAVY; nv.line.fill.background()
    s.shapes.add_picture(str(LOGO_PATH), Inches(0.8), Inches(0.7),
                          width=Inches(3.5), height=Inches(1.2))
    add_textbox(s, Inches(0.8), Inches(2.5), Inches(12), Inches(0.5),
                "知能情報科学セミナーⅠ", size=18, color=NAVY)
    add_textbox(s, Inches(0.8), Inches(3.0), Inches(12), Inches(1.2),
                "モーションキャプチャによる\n歩行・走行・スキップの定量解析",
                size=38, bold=True, color=NAVY)
    add_yellow_underline(s, 5.0, 2.0)
    add_textbox(s, Inches(0.8), Inches(5.15), Inches(12), Inches(0.45),
                "── 急いでいる時はスキップで移動することをお勧めします ──",
                size=15, color=GRAY)
    add_textbox(s, Inches(0.8), Inches(5.85), Inches(12), Inches(0.4),
                "情報学部 情報科学科 3年", size=14)
    add_textbox(s, Inches(0.8), Inches(6.25), Inches(12), Inches(0.5),
                "平松 瑠希", size=22, bold=True, color=NAVY)


def s_intro():
    s = prs.slides.add_slide(BLANK); add_bg(s, WHITE); add_logo(s)
    add_section_badge(s, "01  INTRODUCTION")
    add_title(s, "背景・目的")
    add_yellow_underline(s, 1.75)
    add_textbox(s, Inches(0.5), Inches(2.0), Inches(12.5), Inches(0.6),
                "本レポートは、知能情報科学セミナーⅠで学んだ", size=15)
    add_textbox(s, Inches(0.5), Inches(2.4), Inches(12.5), Inches(0.5),
                "FFT に基づく自己相関 (ACF) ・相互相関 (CCF) の理論を実データに適用する実践演習",
                size=15, bold=True, color=NAVY)
    add_textbox(s, Inches(0.5), Inches(3.3), Inches(12), Inches(0.5),
                "研究の問い", size=20, bold=True, color=NAVY)
    add_bullets(s, Inches(0.5), Inches(3.9), Inches(12), Inches(2.5), [
        "「歩く・走る・スキップ」の動作にはどのようなリズム特性があるか？",
        "同じ動作でも個人差・左右差はどう現れるか？",
        "歩き方の癖・腕振りパターンは ACF でどこまで定量化できるか？",
        "そして、最も移動効率の良い動作は何か？",
    ], size=16)


def s_data():
    s = prs.slides.add_slide(BLANK); add_bg(s, WHITE); add_logo(s)
    add_section_badge(s, "02  DATA")
    add_title(s, "計測データ概要")
    add_yellow_underline(s, 1.75)
    add_bullets(s, Inches(0.5), Inches(2.0), Inches(12), Inches(2), [
        "計測機器: mocopi（ソニー製モーションキャプチャ・IMU センサー 6 個）",
        "フォーマット: BVH (BioVision Hierarchy) / 60 Hz サンプリング",
        "被験者: 三橋青空・長岡翔太・平松瑠希・松田宣久 の 4 名",
        "動作: 歩く・走る・スキップ × 4 名 = 計 12 サンプル",
        "記録項目: root の絶対位置 + 25 関節の回転角（IMU 実測 6・AI 補完 19）",
    ], size=15)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.5), Inches(5.5), Inches(12.333), Inches(1.4))
    box.fill.solid(); box.fill.fore_color.rgb = MUTED_BG
    box.line.color.rgb = NAVY
    add_textbox(s, Inches(0.8), Inches(5.65), Inches(12), Inches(0.4),
                "主要な解析対象信号", size=13, bold=True, color=NAVY)
    add_textbox(s, Inches(0.8), Inches(6.05), Inches(12), Inches(0.8),
                "主：root の鉛直加速度（歩行リズム検出） / 副：左右足首の Xrotation（背屈・底屈）",
                size=12)


def s_method():
    s = prs.slides.add_slide(BLANK); add_bg(s, WHITE); add_logo(s)
    add_section_badge(s, "03  METHOD")
    add_title(s, "解析手法")
    add_yellow_underline(s, 1.75)
    cols = [
        ("① ストライド周期 T（ACF）",
         "ACF(τ) = IFFT(|X(f)|²)",
         "ウィーナー＝ヒンチン定理\n最初の有意ピーク（>0.15）を T と推定"),
        ("② 歩調・歩数",
         "cadence = 1 / T  [Hz]\nsteps = (duration / T) × 2",
         "1 秒あたりのステップ数\n計測時間内の総歩数"),
        ("③ 左右対称性スコア（CCF）",
         "CCF(τ) = IFFT(X*(f) · Y(f))\nScore = 0.6·RMS比\n  + 0.3·位相差 + 0.1·CCFピーク",
         "100 点満点の対称性スコア"),
    ]
    col_w = Inches(4.0)
    for i, (title, formula, desc) in enumerate(cols):
        left = Inches(0.5 + i * 4.27)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  left, Inches(2.1), col_w, Inches(4.7))
        box.fill.solid(); box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = NAVY
        add_textbox(s, left + Inches(0.2), Inches(2.3),
                    col_w - Inches(0.4), Inches(0.6),
                    title, size=13, bold=True, color=NAVY)
        add_textbox(s, left + Inches(0.2), Inches(3.1),
                    col_w - Inches(0.4), Inches(2.4),
                    formula, size=12, color=DARKGRAY, font="Menlo")
        add_textbox(s, left + Inches(0.2), Inches(5.5),
                    col_w - Inches(0.4), Inches(1.3),
                    desc, size=11, color=GRAY)


def s_reading_intro():
    s = prs.slides.add_slide(BLANK); add_bg(s, WHITE); add_logo(s)
    add_section_badge(s, "READING GUIDE")
    add_title(s, "瑠希・歩くを教材に、グラフの読み方を学ぶ")
    add_yellow_underline(s, 1.95)
    add_textbox(s, Inches(0.5), Inches(2.3), Inches(12.5), Inches(2.5),
                "12 サンプル × 各 3 枚 = 36 枚以上のグラフが続く。\n"
                "ここでは「瑠希・歩く」を題材に、グラフを 6 ステップで読み解く方法を示す。\n"
                "ここで読み方を身につければ、後続の 11 サンプルも同じ要領で読める。",
                size=16)
    steps = [
        "STEP 1  腰の上下動を見る ─ 歩行リズムの正体",
        "STEP 2  足首のリズムは腰の「半分」 ─ 2 倍周期の謎",
        "STEP 3  左右の足を重ねる ─ 半周期ずれて動く",
        "STEP 4  ACF でリズムを数値化",
        "STEP 5  CCF で左右対称性を数値化",
        "STEP 6  対称性スコアの内訳 ─ 「何が個性か」を読む",
    ]
    add_bullets(s, Inches(0.5), Inches(4.7), Inches(12.5), Inches(2.5),
                steps, size=13, color=NAVY)


def s_reading_step(num, title, img, point):
    s = prs.slides.add_slide(BLANK); add_bg(s, WHITE); add_logo(s)
    add_section_badge(s, f"READING GUIDE  ─  STEP {num}")
    add_title(s, title, top_in=0.95, size=22)
    add_yellow_underline(s, 1.55)
    add_image_centered(s, img, top_in=1.85, max_width_in=11.5, max_height_in=4.4)
    add_finding(s, "POINT", "ここがポイント", point, top_in=6.4, height_in=0.9)


def s_cwt_intro():
    s = prs.slides.add_slide(BLANK); add_bg(s, WHITE); add_logo(s)
    add_section_badge(s, "★ NEW ANALYSIS ①")
    add_title(s, "連続ウェーブレット変換（CWT）─ 時間軸でリズムの揺らぎを追う", size=22)
    add_yellow_underline(s, 1.65)
    add_textbox(s, Inches(0.5), Inches(2.0), Inches(12.5), Inches(2.5),
                "ACF は「区間全体の平均的リズム」を 1 つの数値に落とすが、時間内のリズムの揺らぎは見えない。\n"
                "CWT は信号を 時間 × 周波数 の 2 次元マップ（ヒートマップ）に展開する。\n"
                "「いつ・どの周波数が強かったか」が色で可視化される。",
                size=14)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.5), Inches(4.5), Inches(12.333), Inches(2.2))
    box.fill.solid(); box.fill.fore_color.rgb = MUTED_BG; box.line.color.rgb = NAVY
    add_textbox(s, Inches(0.8), Inches(4.65), Inches(12), Inches(0.4),
                "本レポートでの設定", size=13, bold=True, color=NAVY)
    add_bullets(s, Inches(0.8), Inches(5.05), Inches(12), Inches(1.7), [
        "母ウェーブレット: Morlet（時間-周波数バランス◎）",
        "入力信号: 腰（root）の鉛直加速度",
        "解析周波数帯: 0.5 〜 5 Hz（歩く 〜 走るをカバー）",
        "4 人分のパワーを揃え正規化 → 平均 → 共通リズム帯を抽出",
    ], size=11)


def s_cwt_motion(motion_ja, img, headline, bullets):
    s = prs.slides.add_slide(BLANK); add_bg(s, WHITE); add_logo(s)
    add_section_badge(s, "★ NEW ANALYSIS ①  CWT")
    add_title(s, f"CWT - {motion_ja}", top_in=0.95, size=22)
    add_yellow_underline(s, 1.55)
    add_textbox(s, Inches(0.5), Inches(1.6), Inches(12.5), Inches(0.5),
                headline, size=14, bold=True, color=NAVY)
    add_image_centered(s, img, top_in=2.1, max_width_in=11.0, max_height_in=4.0)
    add_bullets(s, Inches(0.5), Inches(6.3), Inches(12.5), Inches(1.0),
                bullets, size=11)


def s_cwt_finding():
    s = prs.slides.add_slide(BLANK); add_bg(s, WHITE); add_logo(s)
    add_section_badge(s, "★ NEW ANALYSIS ①  FINDING")
    add_title(s, "スキップは「人類共通リズム」 ─ 誰でも自然にできる動き", size=24)
    add_yellow_underline(s, 1.75)
    add_table(s, ["動作", "4人平均 優位周波数", "個人差 σ", "解釈"],
              [["歩く", "1.90 Hz (114 BPM)", "±0.32 Hz", "個人差 大"],
               ["走る", "3.06 Hz (184 BPM)", "±0.11 Hz", "5 秒後にロックオン"],
               ["スキップ", "2.18 Hz (131 BPM)", "±0.08 Hz", "個人差 最小 ★"]],
              left_in=0.7, top_in=2.3, width_in=12, height_in=2.4)
    add_finding(s, "スキップ推奨の根拠 ④", "スキップは人類共通リズム ─ 習得不要",
                "個人差 σ = 0.08 Hz は 3 動作中最小。4 人とも同じリズム帯（約 2.18 Hz）に収束。"
                "練習なしで誰でも同じリズムでスキップできる ─ 移動手段としての習得コストがゼロ。",
                top_in=5.4)


def s_dtw_intro():
    s = prs.slides.add_slide(BLANK); add_bg(s, WHITE); add_logo(s)
    add_section_badge(s, "★ NEW ANALYSIS ②")
    add_title(s, "動的時間伸縮（DTW）─ 波形そのものの距離を測る", size=24)
    add_yellow_underline(s, 1.75)
    add_textbox(s, Inches(0.5), Inches(2.1), Inches(12.5), Inches(2.5),
                "ACF が「自分自身の周期」、CWT が「時間 × 周波数」を見るのに対し、\n"
                "DTW は「2 つの波形の形そのもの」を比較する。\n"
                "時間軸を伸び縮みさせて最も似た対応関係を探し、距離をスカラ値で出す。\n"
                "「誰と誰がどれくらい似てるか」を 4×4 マトリクスで一望できる。",
                size=14)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.5), Inches(5.0), Inches(12.333), Inches(1.7))
    box.fill.solid(); box.fill.fore_color.rgb = MUTED_BG; box.line.color.rgb = NAVY
    add_textbox(s, Inches(0.8), Inches(5.15), Inches(12), Inches(0.4),
                "本レポートでの設定", size=13, bold=True, color=NAVY)
    add_bullets(s, Inches(0.8), Inches(5.5), Inches(12), Inches(1.2), [
        "各人の左足首 Xrot から 1 ストライド分 → 200 サンプルにリサンプル + 振幅正規化",
        "4 人 × 4 人 = 6 ペアで DTW 距離を計算 → 動作別マトリクス",
    ], size=11)


def s_dtw_image(title, img, bullets):
    s = prs.slides.add_slide(BLANK); add_bg(s, WHITE); add_logo(s)
    add_section_badge(s, "★ NEW ANALYSIS ②  DTW")
    add_title(s, title, top_in=0.95, size=20)
    add_yellow_underline(s, 1.55)
    add_image_centered(s, img, top_in=1.85, max_width_in=11.0, max_height_in=4.3)
    add_bullets(s, Inches(0.5), Inches(6.3), Inches(12.5), Inches(0.9),
                bullets, size=11)


def s_dtw_finding():
    s = prs.slides.add_slide(BLANK); add_bg(s, WHITE); add_logo(s)
    add_section_badge(s, "★ NEW ANALYSIS ②  FINDING")
    add_title(s, "スキップは「波形」も人類共通 ─ 動きの形が決まっている", size=22)
    add_yellow_underline(s, 1.85)
    add_table(s, ["動作", "DTW 距離 平均", "ばらつき σ", "解釈"],
              [["歩く", "30.5", "±14.5", "個人差 中"],
               ["走る", "78.2", "±65.7", "Sora⇔Nori 外れ値あり"],
               ["スキップ", "35.3", "±6.4", "ばらつき 最小 ★"]],
              left_in=0.7, top_in=2.3, width_in=12, height_in=2.4)
    add_finding(s, "スキップ推奨の根拠 ⑤", "スキップは波形も 4 人で揃う ─ 動作の安定性",
                "DTW 距離のばらつき σ = ±6.4 は 3 動作中最小。4 人とも同じ波形をなぞる。"
                "CWT（周波数）と DTW（波形）の 2 つの独立した手法が同じ結論。",
                top_in=5.4)


def s_results_person(person_ja, person_en, color_name, rows, feature):
    s = prs.slides.add_slide(BLANK); add_bg(s, WHITE); add_logo(s)
    add_section_badge(s, "04  RESULTS")
    add_title(s, f"{person_ja}（{person_en}） の全動作データ", top_in=0.95, size=22)
    add_yellow_underline(s, 1.55)
    headers = ["動作", "時間", "T", "歩調", "歩数", "対称性", "速度", "CV"]
    add_table(s, headers, rows,
              left_in=0.5, top_in=2.0, width_in=12.3, height_in=2.2, font_size=12)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.5), Inches(5.0), Inches(12.333), Inches(1.8))
    box.fill.solid(); box.fill.fore_color.rgb = MUTED_BG; box.line.color.rgb = NAVY
    add_textbox(s, Inches(0.8), Inches(5.15), Inches(12), Inches(0.4),
                "特徴", size=14, bold=True, color=NAVY)
    add_textbox(s, Inches(0.8), Inches(5.55), Inches(11.7), Inches(1.2),
                feature, size=12)


def s_4person_compare():
    s = prs.slides.add_slide(BLANK); add_bg(s, WHITE); add_logo(s)
    add_section_badge(s, "04  RESULTS")
    add_title(s, "4 名比較で見える 5 つの特徴", top_in=0.95, size=22)
    add_yellow_underline(s, 1.55)
    bullets = [
        "走るペースは 4 人中 3 人が一致 ─ 青空・翔太・宣久 が 3.16 Hz (190 BPM)。瑠希のみ 171 BPM",
        "歩くペースの個人差は大きい ─ 青空 128 / 宣久 124 / 翔太 116 / 瑠希 112 BPM",
        "対称性最高は青空・走る（93.4 点）",
        "動作の安定性は瑠希が最高 ─ 全動作で CV ≤ 0.017",
        "スキップ速度は走るに近い ─ 全員で 90% 以上、青空は走るを超える",
    ]
    add_bullets(s, Inches(0.5), Inches(2.2), Inches(12.5), Inches(4),
                bullets, size=14)


def s_discovery():
    s = prs.slides.add_slide(BLANK); add_bg(s, WHITE); add_logo(s)
    add_section_badge(s, "05  DISCOVERY")
    add_title(s, "発見① ─ 走るリズムは個人差を超えて収束する", top_in=0.95, size=22)
    add_yellow_underline(s, 1.55)
    big = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(2.0), Inches(2.3), Inches(9.3), Inches(3.5))
    big.fill.solid(); big.fill.fore_color.rgb = NAVY; big.line.fill.background()
    add_textbox(s, Inches(2.0), Inches(2.6), Inches(9.3), Inches(0.5),
                "4 名のうち 3 名の走る歩調が完全一致",
                size=18, color=YELLOW, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(2.0), Inches(3.3), Inches(9.3), Inches(1.5),
                "190 BPM",
                size=90, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(2.0), Inches(5.0), Inches(9.3), Inches(0.4),
                "青空・翔太・宣久 = 3.16 Hz / 瑠希のみ 171 BPM",
                size=14, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.5), Inches(6.2), Inches(12.5), Inches(1.0),
                "170-190 BPM はランニング最適ピッチ域（既存研究と整合）"
                " ─ 身体能力差を超えた人類共通の最適リズムが観測された",
                size=13, color=DARKGRAY, align=PP_ALIGN.CENTER)


def s_extended(title, img, top_text, bullets):
    s = prs.slides.add_slide(BLANK); add_bg(s, WHITE); add_logo(s)
    add_section_badge(s, "06  EXTENDED")
    add_title(s, title, top_in=0.95, size=20)
    add_yellow_underline(s, 1.55)
    add_textbox(s, Inches(0.5), Inches(1.7), Inches(12.5), Inches(0.5),
                top_text, size=13, bold=True, color=NAVY)
    add_image_centered(s, img, top_in=2.2, max_width_in=10.5, max_height_in=4.0)
    add_bullets(s, Inches(0.5), Inches(6.3), Inches(12.5), Inches(0.9),
                bullets, size=11)


def s_gait_characteristics():
    s = prs.slides.add_slide(BLANK); add_bg(s, WHITE); add_logo(s)
    add_section_badge(s, "07  GAIT CHARACTERISTICS")
    add_title(s, "歩き方の特徴 ─ 4 名の比較プロファイル", top_in=0.95, size=22)
    add_yellow_underline(s, 1.55)
    persons = [
        ("三橋 青空", "Sora", "規則的で安定した動作。全動作で CV 極小。スキップが走るより速い特異例"),
        ("長岡 翔太", "Shota", "動作モードを切り替える表現力。歩くで対称性 65.8 点（要分析）"),
        ("平松 瑠希", "Ruki", "全動作で左右非対称度 0.0%、CV ≤ 0.017。4 名で最も安定した動作"),
        ("松田 宣久", "Nori", "全動作で左右非対称度 0.0%。機械的な対称性"),
    ]
    col_w = Inches(3.0)
    for i, (name, en, desc) in enumerate(persons):
        left = Inches(0.5 + i * 3.2)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  left, Inches(2.0), col_w, Inches(4.8))
        box.fill.solid(); box.fill.fore_color.rgb = WHITE; box.line.color.rgb = NAVY
        # Header navy
        hd = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(2.0),
                                 col_w, Inches(0.7))
        hd.fill.solid(); hd.fill.fore_color.rgb = NAVY; hd.line.fill.background()
        add_textbox(s, left + Inches(0.2), Inches(2.1),
                    col_w - Inches(0.4), Inches(0.6),
                    f"{name}", size=14, bold=True, color=WHITE)
        add_textbox(s, left + Inches(0.2), Inches(2.4),
                    col_w - Inches(0.4), Inches(0.3),
                    en, size=10, color=YELLOW)
        add_textbox(s, left + Inches(0.2), Inches(3.0),
                    col_w - Inches(0.4), Inches(3.6),
                    desc, size=12, line_spacing=1.4)


def s_individual(person_ja, motion_ja, img_prefix):
    """個別解析: overview / acf / symmetry を 3 枚並べる"""
    s = prs.slides.add_slide(BLANK); add_bg(s, WHITE); add_logo(s)
    add_section_badge(s, "08  INDIVIDUAL")
    add_title(s, f"個別解析: {person_ja}・{motion_ja}", top_in=0.95, size=22)
    add_yellow_underline(s, 1.55)
    pngs = [
        (f"{img_prefix}_overview.png", "Root 位置・加速度・足首回転角"),
        (f"{img_prefix}_acf.png", "自己相関 ACF"),
        (f"{img_prefix}_symmetry.png", "左右対称性 CCF"),
    ]
    col_w = 3.95
    for i, (fname, cap) in enumerate(pngs):
        path = LG / fname
        if not path.exists():
            continue
        with Image.open(path) as im: iw, ih = im.size
        aspect = iw / ih; max_h = 4.5
        h = min(max_h, col_w / aspect); w = h * aspect
        left = Inches(0.5 + i * 4.27 + (col_w - w) / 2)
        top = Inches(2.0 + (max_h - h) / 2)
        s.shapes.add_picture(str(path), left, top, width=Inches(w), height=Inches(h))
        add_textbox(s, Inches(0.5 + i * 4.27), Inches(6.7), Inches(col_w), Inches(0.4),
                    cap, size=11, color=NAVY, align=PP_ALIGN.CENTER, bold=True)


def s_symmetry_table():
    s = prs.slides.add_slide(BLANK); add_bg(s, WHITE); add_logo(s)
    add_section_badge(s, "09  SYMMETRY")
    add_title(s, "左右対称性スコア一覧", top_in=0.95, size=22)
    add_yellow_underline(s, 1.55)
    add_image_centered(s, LG / "chart_symmetry.png",
                        top_in=2.0, max_width_in=11.0, max_height_in=4.8)
    add_textbox(s, Inches(0.5), Inches(6.9), Inches(12.5), Inches(0.5),
                "瑠希・宣久は全動作で完璧対称（0.0%）、翔太は走る・スキップで完璧、"
                "青空は走る・スキップで対称・歩くで 2 倍関係を検出",
                size=12, color=DARKGRAY, align=PP_ALIGN.CENTER)


def s_speed():
    s = prs.slides.add_slide(BLANK); add_bg(s, WHITE); add_logo(s)
    add_section_badge(s, "10  SPEED ANALYSIS")
    add_title(s, "移動速度の詳細データ ─ 12 サンプル", top_in=0.95, size=22)
    add_yellow_underline(s, 1.55)
    add_image_centered(s, LG / "chart_speed.png",
                        top_in=1.85, max_width_in=11.5, max_height_in=5.2)


def s_key_finding():
    s = prs.slides.add_slide(BLANK); add_bg(s, WHITE); add_logo(s)
    add_section_badge(s, "11  KEY FINDING")
    add_title(s, "発見② ─ スキップは走るの 90% 以上の速度", top_in=0.95, size=22)
    add_yellow_underline(s, 1.55)
    big = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(2.0), Inches(2.3), Inches(9.3), Inches(3.5))
    big.fill.solid(); big.fill.fore_color.rgb = NAVY; big.line.fill.background()
    add_textbox(s, Inches(2.0), Inches(2.6), Inches(9.3), Inches(0.5),
                "4 名の「スキップ ÷ 走る」速度比 平均",
                size=18, color=YELLOW, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(2.0), Inches(3.3), Inches(9.3), Inches(1.5),
                "94.4%", size=90, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(2.0), Inches(5.0), Inches(9.3), Inches(0.4),
                "青空 100.2% / 翔太 98.8% / 瑠希 89.6% / 宣久 89.3%",
                size=14, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.5), Inches(6.2), Inches(12.5), Inches(1.0),
                "全員でスキップは走るの約 90% 以上の速度。青空に至っては走るを超える",
                size=13, color=DARKGRAY, align=PP_ALIGN.CENTER, bold=True)


def s_efficiency():
    s = prs.slides.add_slide(BLANK); add_bg(s, WHITE); add_logo(s)
    add_section_badge(s, "12  EFFICIENCY")
    add_title(s, "リズム vs 速度の関係 ─ 4 名比較", top_in=0.95, size=22)
    add_yellow_underline(s, 1.55)
    add_image_centered(s, LG / "chart_efficiency.png",
                        top_in=1.85, max_width_in=11.0, max_height_in=5.2)


def s_discussion():
    s = prs.slides.add_slide(BLANK); add_bg(s, WHITE); add_logo(s)
    add_section_badge(s, "13  DISCUSSION")
    add_title(s, "考察", top_in=0.95, size=24)
    add_yellow_underline(s, 1.55)
    points = [
        ("① 走るリズムの普遍性",
         "4 人中 3 人が走る 3.16 Hz (190 BPM) 一致。瑠希のみ 171 BPM。"
         "身体能力差を超えた人間共通の最適ピッチを示唆"),
        ("② 歩き方の個人差",
         "左右 ACF・腕足周期比で歩く・スキップに個人差が出る。"
         "走るは身体力学的制約が強く 4 名で似通る"),
        ("③ 4 人それぞれの動作スタイル",
         "青空: 規則的・安定 / 翔太: モード切替 / "
         "瑠希: 左右対称 0%・CV 最小 / 宣久: 機械的対称性"),
        ("④ スキップの優位性",
         "経路速度 走るの 90% 以上 / ピッチは走るの約半分 / "
         "対称性スコア 77-93 点 / 左右非対称度 4 人中 3 人で 0.0%"),
        ("⑤ 今後の課題",
         "翔太の歩く対称性低下の原因 / サンプル数増加 / エネルギー消費との対応"),
    ]
    y = 1.8
    for title, body in points:
        add_textbox(s, Inches(0.5), Inches(y), Inches(12.5), Inches(0.35),
                    title, size=13, bold=True, color=NAVY)
        add_textbox(s, Inches(0.5), Inches(y + 0.35), Inches(12.5), Inches(0.65),
                    body, size=11)
        y += 1.05


def s_tempo():
    s = prs.slides.add_slide(BLANK); add_bg(s, WHITE); add_logo(s)
    add_section_badge(s, "14  BONUS  ─  TEMPO")
    add_title(s, "スキップのBPMは、心地よいポップスのテンポと一致する",
              top_in=0.95, size=20)
    add_yellow_underline(s, 1.65)
    add_textbox(s, Inches(0.5), Inches(1.95), Inches(12.5), Inches(0.5),
                "4 名のスキップ BPM 範囲: 90 - 103",
                size=15, bold=True, color=NAVY)
    songs = [
        ("aimyon", "マリーゴールド", "104 BPM"),
        ("YOASOBI", "夜に駆ける", "130 BPM（疾走）"),
        ("Mr.Children", "innocent world", "108 BPM"),
        ("Official髭男dism", "Pretender", "92 BPM"),
        ("ヨルシカ", "夜行", "100 BPM"),
        ("優里", "ドライフラワー", "78 BPM（緩やか）"),
    ]
    col_w = Inches(4.0)
    for i, (artist, song, bpm) in enumerate(songs):
        col = i % 3; row = i // 3
        left = Inches(0.5 + col * 4.27)
        top = Inches(2.7 + row * 1.7)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  left, top, col_w, Inches(1.4))
        box.fill.solid(); box.fill.fore_color.rgb = MUTED_BG; box.line.color.rgb = NAVY
        add_textbox(s, left + Inches(0.2), top + Inches(0.1),
                    col_w - Inches(0.4), Inches(0.4),
                    artist, size=11, color=GRAY)
        add_textbox(s, left + Inches(0.2), top + Inches(0.45),
                    col_w - Inches(0.4), Inches(0.5),
                    song, size=15, bold=True, color=NAVY)
        add_textbox(s, left + Inches(0.2), top + Inches(0.95),
                    col_w - Inches(0.4), Inches(0.45),
                    bpm, size=12, color=YELLOW, bold=True)
    add_textbox(s, Inches(0.5), Inches(6.8), Inches(12.5), Inches(0.5),
                "→ スキップのリズムは「気持ちよく聴ける音楽のテンポ」と物理的に同期する",
                size=13, color=DARKGRAY, align=PP_ALIGN.CENTER, bold=True)


def s_conclusion_main():
    s = prs.slides.add_slide(BLANK); add_bg(s, NAVY)
    # Yellow bar
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.25))
    bar.fill.solid(); bar.fill.fore_color.rgb = YELLOW; bar.line.fill.background()
    add_logo(s)
    add_section_badge(s, "15  CONCLUSION")
    add_textbox(s, Inches(0.5), Inches(1.5), Inches(12.5), Inches(0.5),
                "FINAL CONCLUSION", size=14, bold=True, color=YELLOW,
                align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.5), Inches(2.5), Inches(12.5), Inches(3.5),
                "急いでいる時は\nスキップで移動することを\nお勧めします",
                size=54, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, line_spacing=1.3)


def s_conclusion_reasons():
    s = prs.slides.add_slide(BLANK); add_bg(s, WHITE); add_logo(s)
    add_section_badge(s, "15  CONCLUSION")
    add_title(s, "なぜスキップを推奨するのか ─ 5 つの根拠", top_in=0.95, size=22)
    add_yellow_underline(s, 1.55)
    cards = [
        ("①", "速度", "走るの 90% 以上維持\n青空はスキップが速い\n4名平均 94.4%"),
        ("②", "リズム", "ピッチは走るの半分\n90-103 BPM\nマリーゴールドと同テンポ ♪"),
        ("③", "対称性", "4人中3人が左右非対称度\n0.0%\n左右バランス◎"),
        ("④", "習得不要", "CWT 個人差 σ\n= 0.08 Hz（最小）\n誰でも同じリズム"),
        ("⑤", "動作の安定", "DTW 波形ばらつき σ\n= ±6.4（最小）\n4 人とも同じ形をなぞる"),
    ]
    col_w = Inches(2.4)
    for i, (num, title, desc) in enumerate(cards):
        left = Inches(0.5 + i * 2.55)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  left, Inches(2.0), col_w, Inches(4.6))
        box.fill.solid(); box.fill.fore_color.rgb = WHITE; box.line.color.rgb = NAVY
        # Top yellow band
        hd = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(2.0),
                                 col_w, Inches(0.8))
        hd.fill.solid(); hd.fill.fore_color.rgb = NAVY; hd.line.fill.background()
        add_textbox(s, left, Inches(2.05), col_w, Inches(0.7),
                    num, size=38, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
        add_textbox(s, left + Inches(0.1), Inches(3.0),
                    col_w - Inches(0.2), Inches(0.5),
                    title, size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_textbox(s, left + Inches(0.15), Inches(3.6),
                    col_w - Inches(0.3), Inches(2.9),
                    desc, size=11, align=PP_ALIGN.CENTER, line_spacing=1.4)
    add_textbox(s, Inches(0.5), Inches(6.9), Inches(12.5), Inches(0.5),
                "走るのがしんどい・歩くのは遅い、という場面において、"
                "スキップは合理的かつ楽しい選択肢である",
                size=13, color=DARKGRAY, align=PP_ALIGN.CENTER, bold=True)


def s_references():
    s = prs.slides.add_slide(BLANK); add_bg(s, WHITE); add_logo(s)
    add_section_badge(s, "REFERENCES")
    add_title(s, "参考文献", size=28)
    add_yellow_underline(s, 1.85)
    refs = [
        ("分析", "[1] Cooley & Tukey (1965). Mathematics of Computation, 19(90), 297–301."),
        ("分析", "[2] Sakoe & Chiba (1978). IEEE TASSP, 26(1), 43–49."),
        ("分析", "[3] Mallat (1989). IEEE PAMI, 11(7), 674–693."),
        ("分析", "[4] Torrence & Compo (1998). BAMS, 79(1), 61–78."),
        ("歩行", "[5] Cavagna et al. (1977). Am J Physiol, 233(5), R243–R261."),
        ("歩行", "[6] Alexander (1989). Physiol Rev, 69(4), 1199–1227."),
        ("歩行", "[7] Hreljac (1995). J Biomech, 28(6), 669–677."),
        ("歩行", "[8] Minetti (1998). Proc R Soc B, 265(1402), 1227–1235."),
        ("歩行", "[9] Pavei et al. (2015). J Appl Physiol, 119(1), 93–100."),
    ]
    tb = s.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(12.5), Inches(5))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (cat, ref) in enumerate(refs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(5)
        r1 = p.add_run(); r1.text = f"[{cat}] "
        r1.font.size = Pt(11); r1.font.bold = True
        r1.font.color.rgb = NAVY if cat == "分析" else RGBColor(0x2C, 0x8C, 0x4F)
        r1.font.name = "Hiragino Sans"
        r2 = p.add_run(); r2.text = ref
        r2.font.size = Pt(11); r2.font.color.rgb = DARKGRAY
        r2.font.name = "Hiragino Sans"


def s_closing():
    s = prs.slides.add_slide(BLANK); add_bg(s, NAVY)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.25))
    bar.fill.solid(); bar.fill.fore_color.rgb = YELLOW; bar.line.fill.background()
    s.shapes.add_picture(str(LOGO_PATH), Inches(5.4), Inches(1.0),
                          width=Inches(2.5), height=Inches(0.85))
    add_textbox(s, Inches(0.5), Inches(3.0), Inches(12.5), Inches(1.5),
                "ご清聴ありがとうございました",
                size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.5), Inches(4.7), Inches(12.5), Inches(0.5),
                "知能情報科学セミナーⅠ ─ BVH動作解析レポート",
                size=18, color=YELLOW, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.5), Inches(5.4), Inches(12.5), Inches(0.5),
                "情報学部 情報科学科  ─  平松 瑠希",
                size=16, color=WHITE, align=PP_ALIGN.CENTER)


def add_footer(slide, number, total):
    add_textbox(slide, Inches(0.5), SLIDE_H - Inches(0.4),
                Inches(8), Inches(0.3),
                "知能情報科学セミナーⅠ ── 平松瑠希 ── 工学院大学",
                size=8, color=GRAY)
    add_textbox(slide, SLIDE_W - Inches(1.5), SLIDE_H - Inches(0.4),
                Inches(1.0), Inches(0.3),
                f"{number} / {total}",
                size=8, color=GRAY, align=PP_ALIGN.RIGHT)


# ===== ビルド =====
def main():
    s_title()
    s_intro()
    s_data()
    s_method()

    s_reading_intro()
    s_reading_step(1, "腰の上下動を見る ─ 歩行リズムの正体",
                   RG / "img_01_root_position.png",
                   "ピーク間隔の平均 ≈ 0.56 s ＝ 半ストライド（T/2）。腰は 1 ストライドの間に 2 回上下する。")
    s_reading_step(2, "足首のリズムは腰の「半分」 ─ 2 倍周期の謎",
                   RG / "img_02_root_vs_foot.png",
                   "左足は 1 ストライドに 1 回しか動かないが、腰は 2 回上下する。足首周期 1.05 s ≈ 腰周期 0.53 s × 2。")
    s_reading_step(3, "左右の足を重ねる ─ 半周期ずれて動く",
                   RG / "img_03_left_right.png",
                   "左右ピーク間 ≈ 0.58 s ＝ T/2。これが「交互ステップ」の数学的定義。")
    s_reading_step(4, "ACF でリズムを数値化",
                   RG / "img_04_acf.png",
                   "τ = 0.53 s でピーク → ストライド周期 T。歩調 1.87 Hz ＝ 112 BPM。2 秒以上ピークが残る = 規則的で安定。")
    s_reading_step(5, "CCF で左右対称性を数値化",
                   RG / "img_05_ccf.png",
                   "ピーク位置 -0.517 s ≈ 理想 -T/2 = -0.533 s。差わずか 0.016 s ─ 完璧な交互ステップ。")
    s_reading_step(6, "対称性スコアの内訳 ─ 「何が個性か」を読む",
                   RG / "img_06_score.png",
                   "84.3 点。タイミングはほぼ満点、RMS 比 0.81 で 19% の左右振幅差 → これが歩き方の癖。")

    s_cwt_intro()
    s_cwt_motion("歩く", CWT / "walk_avg.png",
                 "1.90 Hz（114 BPM）、個人差 σ = 0.32 Hz ─ 最も個性が出る",
                 ["1.5〜2.2 Hz の幅で 4 人がばらつく",
                  "歩行は日常動作ゆえに個性が許容される"])
    s_cwt_motion("走る", CWT / "run_avg.png",
                 "3.06 Hz（184 BPM）、個人差 σ = 0.11 Hz ─ 5 秒後にロックオン",
                 ["0〜5 秒は加速期、5 秒以降に 3 Hz の明るい横帯",
                  "4 人とも同タイミングで同周波数に収束"])
    s_cwt_motion("スキップ", CWT / "skip_avg.png",
                 "2.18 Hz、個人差 σ = 0.08 Hz ─ 全動作中最小",
                 ["複数の周波数帯（基本リズム + 高調波）が共存",
                  "「ステップ + ホップ」の複合動作だが優位周波数は最も揃う"])
    s_cwt_finding()

    s_dtw_intro()
    s_dtw_image("4 人ペアワイズ DTW 距離マトリクス",
                DTW / "distance_matrix.png",
                ["色が薄い（黄）ほど波形が似ている、濃い（赤）ほど違う",
                 "走るの Sora ⇔ Nori だけ 216.3 と異常値 ─ 他は 23〜82"])
    s_dtw_image("動作ごとの平均距離 ─ 「ばらつき」で比較",
                DTW / "distance_summary.png",
                ["歩く 30.5 ± 14.5 / 走る 78.2 ± 65.7 / スキップ 35.3 ± 6.4",
                 "標準偏差の小ささ = 4 人がどれだけ揃っているか ─ スキップが最小"])
    s_dtw_image("DTW の中身 ─ Ruki ⇔ Nori（歩く）の例",
                DTW / "warping_example.png",
                ["左：2 人の波形と DTW の対応付け",
                 "右：累積コスト行列。赤線が最適パス"])
    s_dtw_finding()

    # 04 RESULTS：個人別
    s_results_person("三橋 青空", "Sora", "sora",
        [["歩く", "13.6 s", "0.467 s", "2.14 Hz (128 BPM)", "58", "88.2", "6.43 km/h", "0.018"],
         ["走る", "20.2 s", "0.317 s", "3.16 Hz (190 BPM)", "127", "93.4", "11.63 km/h", "0.026"],
         ["スキップ", "19.1 s", "0.583 s", "1.71 Hz (103 BPM)", "65", "92.6", "11.65 km/h ★", "0.737"]],
        "全動作で対称性スコア 88 点以上。歩く・走るの CV が極めて低い → 非常に安定したリズミカルな動作。スキップが走るより僅かに速い。")
    s_results_person("長岡 翔太", "Shota", "shota",
        [["歩く", "53.0 s", "0.517 s", "1.94 Hz (116 BPM)", "205", "65.8", "1.97 km/h", "0.015"],
         ["走る", "29.3 s", "0.317 s", "3.16 Hz (190 BPM)", "185", "85.0", "10.37 km/h", "0.663"],
         ["スキップ", "28.8 s", "0.617 s", "1.62 Hz (97 BPM)", "93", "90.0", "10.25 km/h", "1.894"]],
        "走るリズム 190 BPM で青空と一致。歩く対称性のみ低スコア。スキップ CV=1.89 と動作のばらつきが大きい。")
    s_results_person("平松 瑠希", "Ruki", "ruki",
        [["歩く", "17.2 s", "0.533 s", "1.87 Hz (112 BPM)", "64", "85.7", "7.21 km/h", "0.016"],
         ["走る", "25.3 s", "0.350 s", "2.86 Hz (171 BPM)", "144", "79.7", "10.28 km/h", "0.017"],
         ["スキップ", "28.3 s", "0.667 s", "1.50 Hz (90 BPM)", "85", "77.1", "9.21 km/h", "0.012"]],
        "走るは 171 BPM で 4 名中唯一の例外。全動作で CV ≤ 0.017 と極めて安定。歩くで 7.21 km/h は 4 名中最速。")
    s_results_person("松田 宣久", "Nori", "nori",
        [["歩く", "19.0 s", "0.483 s", "2.07 Hz (124 BPM)", "78", "83.4", "4.91 km/h", "0.526"],
         ["走る", "27.4 s", "0.317 s", "3.16 Hz (190 BPM)", "172", "88.7", "9.45 km/h", "1.468"],
         ["スキップ", "30.5 s", "0.633 s", "1.58 Hz (95 BPM)", "96", "84.3", "8.44 km/h", "0.334"]],
        "走るリズム 190 BPM で青空・翔太と完全一致。拡張解析で全動作の左右非対称度が 0.00% ─ 機械的な左右対称性。")
    s_4person_compare()

    s_discovery()

    # 06 EXTENDED
    s_extended("拡張解析① ─ 左右足首 ACF による左右差の詳細",
               LG / "ext_lr_foot_acf.png",
               "左右の足首 Xrot 各々の ACF から周期非対称度を算出",
               ["瑠希・宣久は全動作で左右非対称度 0.0%（完璧）",
                "青空の歩くで 66.7%（左右 2 倍関係 ─ 片足/ステップの可能性）"])
    s_extended("拡張解析② ─ 腕振りパターン（腕⇔足 相互相関）",
               LG / "ext_arm_leg_ccf.png",
               "腕足周期比から、腕振りが足のリズムと同期しているかを確認",
               ["同期型（青空・瑠希・宣久）─ 歩く・走るで腕足が 1:1 同期、スキップで急変",
                "翔太は歩くで腕振り同期突出"])
    s_extended("拡張解析③ ─ ACF ピーク減衰率（動作の規則性）",
               LG / "ext_acf_decay.png",
               "2T/T, 3T/T ピーク比 = どれだけ周期が安定しているか",
               ["スキップは「跳ねる」動作で減衰が大きい",
                "走るは減衰が小さい = 周期が長期的に安定"])
    s_extended("拡張解析④ ─ 動作プロファイル（6 次元レーダー）",
               LG / "ext_radar.png",
               "歩調・対称性・左右ACF・腕足周期比・速度・CV を 1 枚に集約",
               ["瑠希・宣久は全動作で左右対称軸が満点",
                "青空は走るがバランス取れ、翔太は歩くで腕振り同期突出"])

    s_gait_characteristics()

    # 08 INDIVIDUAL：12 サンプル
    individuals = [
        ("青空", "歩く", "Sora_walking"),
        ("青空", "走る", "Sora_running"),
        ("青空", "スキップ", "Sora_skip"),
        ("翔太", "歩く", "Shota_walking"),
        ("翔太", "走る", "Sho_run"),
        ("翔太", "スキップ", "Shota_skip"),
        ("瑠希", "歩く", "Ruki_walk"),
        ("瑠希", "走る", "Ruki_run"),
        ("瑠希", "スキップ", "Ruki_skip"),
        ("宣久", "歩く", "Nori_walk"),
        ("宣久", "走る", "Nori_run"),
        ("宣久", "スキップ", "Nori_skip"),
    ]
    for person, motion, prefix in individuals:
        s_individual(person, motion, prefix)

    s_symmetry_table()
    s_speed()
    s_key_finding()
    s_efficiency()
    s_discussion()
    s_tempo()

    s_conclusion_main()
    s_conclusion_reasons()
    s_references()
    s_closing()

    # ページ番号フッタ
    total = len(prs.slides)
    for i, sl in enumerate(prs.slides, start=1):
        add_footer(sl, i, total)

    prs.save(str(OUT))
    print(f"✅ PPTX 生成完了: {OUT}")
    print(f"   スライド数: {total}")


if __name__ == "__main__":
    main()
